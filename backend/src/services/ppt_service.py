import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO

logger = logging.getLogger(__name__)

# Luxury Palette Tokens
COLOR_PRIMARY = RGBColor(26, 26, 46)     # Deep Navy #1a1a2e
COLOR_ACCENT = RGBColor(196, 30, 58)     # Crimson Accent #c41e3a
COLOR_GOLD = RGBColor(212, 175, 55)      # Luxury Gold #d4af37
COLOR_TEXT_DARK = RGBColor(30, 41, 59)   # Slate 800
COLOR_TEXT_MUTED = RGBColor(100, 116, 139) # Slate 500
COLOR_SURFACE = RGBColor(248, 250, 252)  # Slate 50
COLOR_WHITE = RGBColor(255, 255, 255)

def add_header_banner(slide, prs, title_text: str, category_text: str = "VOYANTA LUXURY PORTFOLIO"):
    """Adds a magazine-style top header with crimson accent strip and stylish title."""
    # Top Crimson bar
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.15))
    strip.fill.solid()
    strip.fill.fore_color.rgb = COLOR_ACCENT
    strip.line.fill.background()

    # Category Subhead
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.4))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = category_text.upper()
    p_sub.font.size = Pt(11)
    p_sub.font.bold = True
    p_sub.font.color.rgb = COLOR_GOLD
    p_sub.font.name = "Arial"

    # Main Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(32)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_PRIMARY
    p_title.font.name = "Georgia"

def generate_ppt_from_proposal(proposal: dict, items: list) -> bytes:
    """
    Generates an executive widescreen 16:9 magazine-quality PowerPoint presentation
    with custom styling, color palettes, and structured layout slides.
    """
    try:
        prs = Presentation()
        # Set widescreen 16:9 dimensions (13.33 x 7.5 inches)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6] # Blank slide layout
        
        proposal_name = proposal.get("name") or proposal.get("title") or "Bespoke Travel Journey"
        client_name = proposal.get("client_name") or proposal.get("client") or "Valued Client"
        dest = proposal.get("destination") or "Selected Destinations"
        days = proposal.get("duration_days") or proposal.get("days") or 7
        overview = proposal.get("overview") or proposal.get("description") or f"A curated luxury itinerary exploring {dest} with VIP accommodations and private tours."

        # ─── Slide 1: Cinematic Title Cover ──────────────────────────────────
        slide1 = prs.slides.add_slide(blank_layout)
        
        # Dark luxury background fill
        bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_PRIMARY
        bg.line.fill.background()

        # Left Gold Accent bar
        bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), prs.slide_height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_GOLD
        bar.line.fill.background()

        # Title Box
        t_box = slide1.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.5), Inches(2.0))
        tf = t_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = proposal_name
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = "Georgia"

        # Subtitle Box
        s_box = slide1.shapes.add_textbox(Inches(1.5), Inches(4.3), Inches(10.5), Inches(1.5))
        stf = s_box.text_frame
        stf.word_wrap = True
        sp1 = stf.paragraphs[0]
        sp1.text = f"DESTINATION: {dest.upper()}   |   DURATION: {days} DAYS"
        sp1.font.size = Pt(14)
        sp1.font.bold = True
        sp1.font.color.rgb = COLOR_GOLD
        sp1.font.name = "Arial"

        sp2 = stf.add_paragraph()
        sp2.text = f"Prepared Exclusively for: {client_name}"
        sp2.font.size = Pt(22)
        sp2.font.italic = True
        sp2.font.color.rgb = COLOR_SURFACE
        sp2.font.name = "Georgia"

        # ─── Slide 2: Executive Overview & Highlights ────────────────────────
        slide2 = prs.slides.add_slide(blank_layout)
        add_header_banner(slide2, prs, "Executive Summary & Journey Overview", f"DESTINATION FOCUS: {dest}")

        # Overview Card
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.8))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_SURFACE
        card.line.color.rgb = COLOR_GOLD
        card.line.width = Pt(1.5)
        
        ctf = card.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cp.text = "THE EXPERIENCE"
        cp.font.size = Pt(12)
        cp.font.bold = True
        cp.font.color.rgb = COLOR_ACCENT
        cp.font.name = "Arial"

        cp2 = ctf.add_paragraph()
        cp2.text = overview
        cp2.font.size = Pt(16)
        cp2.font.color.rgb = COLOR_TEXT_DARK
        cp2.font.name = "Georgia"

        # Highlights List
        h_box = slide2.shapes.add_textbox(Inches(0.8), Inches(3.9), Inches(11.7), Inches(3.0))
        htf = h_box.text_frame
        htf.word_wrap = True
        hp = htf.paragraphs[0]
        hp.text = "CURATED HIGHLIGHTS & VIP INCLUSIONS:"
        hp.font.size = Pt(16)
        hp.font.bold = True
        hp.font.color.rgb = COLOR_PRIMARY
        hp.font.name = "Georgia"

        highlights = proposal.get("highlights") or [
            "Private chauffeured airport transfers and VIP meet & greet assistance",
            "Hand-selected 5-star accommodations with panoramic views and premier hospitality",
            "Exclusive private guided cultural and historical tours with local experts",
            "Dedicated 24/7 concierge support throughout your entire stay"
        ]
        for h_text in highlights[:5]:
            hp_item = htf.add_paragraph()
            hp_item.text = f"✦  {h_text}"
            hp_item.font.size = Pt(15)
            hp_item.font.color.rgb = COLOR_TEXT_DARK
            hp_item.font.name = "Arial"

        # ─── Slide 3: Daily Itinerary Spreads ────────────────────────────────
        itinerary_days = proposal.get("itinerary") or proposal.get("days_list") or []
        if isinstance(itinerary_days, list) and len(itinerary_days) > 0:
            for idx, d in enumerate(itinerary_days[:8]): # Cap at 8 slides for clean export
                if not isinstance(d, dict): continue
                slide_day = prs.slides.add_slide(blank_layout)
                day_num = d.get("day", idx + 1)
                day_title = d.get("title") or f"Day {day_num} Discovery"
                day_desc = d.get("description") or "Spend the day exploring landmarks and enjoying private hospitality."
                
                add_header_banner(slide_day, prs, f"Day 0{day_num}: {day_title}", "DAILY CHRONICLE")

                # Day Content Box
                box = slide_day.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5))
                box.fill.solid()
                box.fill.fore_color.rgb = COLOR_SURFACE
                box.line.color.rgb = COLOR_TEXT_MUTED
                
                dtf = box.text_frame
                dtf.word_wrap = True
                dp1 = dtf.paragraphs[0]
                dp1.text = "SCHEDULING & ACTIVITIES"
                dp1.font.size = Pt(12)
                dp1.font.bold = True
                dp1.font.color.rgb = COLOR_ACCENT

                dp2 = dtf.add_paragraph()
                dp2.text = day_desc
                dp2.font.size = Pt(18)
                dp2.font.color.rgb = COLOR_TEXT_DARK
                dp2.font.name = "Georgia"

                activities = d.get("activities") or []
                if isinstance(activities, list) and len(activities) > 0:
                    dp3 = dtf.add_paragraph()
                    dp3.text = "\nFeatured Experiences: " + " · ".join([str(a) for a in activities])
                    dp3.font.size = Pt(14)
                    dp3.font.italic = True
                    dp3.font.color.rgb = COLOR_PRIMARY

        # ─── Slide 4: Selected Accommodations & Services ─────────────────────
        grouped = {}
        for it in items:
            if not isinstance(it, dict): continue
            k = (it.get("kind") or "service").lower()
            grouped.setdefault(k, []).append(it)

        for kind, group_items in grouped.items():
            slide_serv = prs.slides.add_slide(blank_layout)
            add_header_banner(slide_serv, prs, f"Featured {kind.capitalize()}s & Reservations", "INVENTORY DETAILS")
            
            top_offset = 1.8
            for idx, item in enumerate(group_items[:4]): # Show up to 4 items per slide
                label = item.get("label") or item.get("name") or f"{kind.capitalize()} Item"
                desc = item.get("desc") or item.get("description") or item.get("notes") or ""
                price = (float(item.get("unit_price") or 0)) * (float(item.get("qty") or 1))

                s_box = slide_serv.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top_offset), Inches(11.7), Inches(1.1))
                s_box.fill.solid()
                s_box.fill.fore_color.rgb = COLOR_WHITE if idx % 2 == 0 else COLOR_SURFACE
                s_box.line.color.rgb = COLOR_GOLD if idx == 0 else COLOR_TEXT_MUTED
                
                stf = s_box.text_frame
                stf.word_wrap = True
                sp = stf.paragraphs[0]
                sp.text = f"✦  {label.upper()}" + (f"  —  ₹{price:,.0f}" if price > 0 else "")
                sp.font.size = Pt(16)
                sp.font.bold = True
                sp.font.color.rgb = COLOR_PRIMARY

                if desc:
                    sp_desc = stf.add_paragraph()
                    sp_desc.text = f"    {desc}"
                    sp_desc.font.size = Pt(13)
                    sp_desc.font.color.rgb = COLOR_TEXT_MUTED

                top_offset += 1.3

        # ─── Slide 5: Sign-Off & Agency Footer ───────────────────────────────
        slide_end = prs.slides.add_slide(blank_layout)
        bg_end = slide_end.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        bg_end.fill.solid()
        bg_end.fill.fore_color.rgb = COLOR_PRIMARY
        bg_end.line.fill.background()

        e_box = slide_end.shapes.add_textbox(Inches(2.0), Inches(2.2), Inches(9.3), Inches(3.0))
        etf = e_box.text_frame
        etf.word_wrap = True
        ep1 = etf.paragraphs[0]
        ep1.text = "Thank You for Your Trust"
        ep1.font.size = Pt(44)
        ep1.font.bold = True
        ep1.font.color.rgb = COLOR_GOLD
        ep1.font.name = "Georgia"
        ep1.alignment = PP_ALIGN.CENTER

        ep2 = etf.add_paragraph()
        ep2.text = "\nWe invite you to approve this itinerary or connect with your travel designer for any bespoke modifications."
        ep2.font.size = Pt(18)
        ep2.font.color.rgb = COLOR_WHITE
        ep2.font.name = "Arial"
        ep2.alignment = PP_ALIGN.CENTER

        branding = proposal.get("preferences", {}).get("branding", {})
        agency = branding.get("agency_name") or "Voyanta Luxury Travel"
        email = branding.get("contact_email") or "concierge@voyantatravel.com"
        phone = branding.get("contact_phone") or "+1 (800) 555-8692"

        ep3 = etf.add_paragraph()
        ep3.text = f"\n\n{agency}   |   {email}   |   {phone}"
        ep3.font.size = Pt(14)
        ep3.font.bold = True
        ep3.font.color.rgb = COLOR_GOLD
        ep3.alignment = PP_ALIGN.CENTER

        # Save to buffer
        ppt_stream = BytesIO()
        prs.save(ppt_stream)
        ppt_stream.seek(0)
        
        return ppt_stream.getvalue()
        
    except Exception as e:
        logger.exception("PPT generation failed")
        raise e
