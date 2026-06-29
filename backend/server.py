from fastapi import FastAPI, APIRouter, Request, HTTPException
from contextlib import asynccontextmanager
from fastapi.responses import Response
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
import os
import logging
import httpx
from pathlib import Path
from pydantic import BaseModel, Field, ConfigDict
from typing import List, Optional, Any
import uuid
from datetime import datetime, timezone
from fastapi import Depends, Security
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from tenacity import retry, stop_after_attempt, wait_exponential


ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

SUPABASE_URL = os.environ.get('VITE_SUPABASE_URL') or os.environ.get('SUPABASE_URL')
SUPABASE_KEY = os.environ.get('VITE_SUPABASE_ANON_KEY') or os.environ.get('SUPABASE_KEY')

@asynccontextmanager
async def lifespan(app: FastAPI):
    yield

# Create the main app without a prefix
app = FastAPI(lifespan=lifespan)

# Create a router with the /api prefix
api_router = APIRouter(prefix="/api")


# Define Models
class StatusCheck(BaseModel):
    model_config = ConfigDict(extra="ignore")  # Ignore MongoDB's _id field
    
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    client_name: str
    timestamp: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

class StatusCheckCreate(BaseModel):
    client_name: str

class ParseItineraryInput(BaseModel):
    text: str

class PDFGenerateRequest(BaseModel):
    html: str
    name: Optional[str] = "proposal"

class PPTGenerateRequest(BaseModel):
    proposal: dict = Field(default_factory=dict)
    items: list = Field(default_factory=list)

security = HTTPBearer()

async def verify_token(credentials: HTTPAuthorizationCredentials = Security(security)):
    token = credentials.credentials
    if not SUPABASE_URL or not SUPABASE_KEY:
        raise HTTPException(status_code=500, detail="Supabase not configured")
        
    headers = {
        "apikey": SUPABASE_KEY,
        "Authorization": f"Bearer {token}"
    }
    try:
        async with httpx.AsyncClient(timeout=5.0) as client:
            res = await client.get(f"{SUPABASE_URL}/auth/v1/user", headers=headers)
            if res.status_code != 200:
                logger.error(f"Token verification failed: {res.text}")
                raise HTTPException(status_code=401, detail="Invalid auth token")
            return res.json()
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to verify token: {e}")
        raise HTTPException(status_code=401, detail="Invalid auth token")

# Add your routes to the router instead of directly to app
@api_router.get("/")
async def root():
    return {"message": "Hello World"}

@api_router.get("/health")
async def health():
    return {
        "status": "ok",
        "version": "1.0.0",
        "timestamp": datetime.now(timezone.utc).isoformat()
    }

@api_router.post("/status", response_model=StatusCheck)
async def create_status_check(input: StatusCheckCreate):
    status_dict = input.model_dump()
    status_obj = StatusCheck(**status_dict)
    
    if not SUPABASE_URL or not SUPABASE_KEY:
        return status_obj
        
    headers = {
        "apikey": SUPABASE_KEY,
        "Authorization": f"Bearer {SUPABASE_KEY}",
        "Content-Type": "application/json",
        "Prefer": "return=minimal"
    }
    doc = status_obj.model_dump()
    doc['timestamp'] = doc['timestamp'].isoformat()
    
    try:
        async with httpx.AsyncClient(timeout=5.0) as client:
            await client.post(f"{SUPABASE_URL}/rest/v1/status_checks", json=doc, headers=headers)
    except Exception as e:
        logger.error(f"Failed to record status check to Supabase: {e}")
        
    return status_obj

@api_router.get("/status", response_model=List[StatusCheck])
async def get_status_checks():
    if not SUPABASE_URL or not SUPABASE_KEY:
        return []
        
    headers = {
        "apikey": SUPABASE_KEY,
        "Authorization": f"Bearer {SUPABASE_KEY}"
    }
    
    try:
        async with httpx.AsyncClient(timeout=5.0) as client:
            res = await client.get(f"{SUPABASE_URL}/rest/v1/status_checks?select=*&order=timestamp.desc", headers=headers)
            if res.status_code == 200:
                data = res.json()
                for check in data:
                    if isinstance(check.get('timestamp'), str):
                        check['timestamp'] = datetime.fromisoformat(check['timestamp'])
                return data
    except Exception as e:
        logger.error(f"Failed to fetch status checks from Supabase: {e}")
        
    return []


@retry(stop=stop_after_attempt(3), wait=wait_exponential(multiplier=1, min=2, max=10))
async def call_openai_with_retry(payload: dict, headers: dict):
    async with httpx.AsyncClient(timeout=60.0) as client:
        r = await client.post("https://api.openai.com/v1/chat/completions", json=payload, headers=headers)
        if r.status_code == 429 or r.status_code >= 500:
            logger.warning(f"OpenAI error {r.status_code}, retrying...")
            raise Exception(f"OpenAI error {r.status_code}: {r.text}")
        if r.status_code != 200:
            raise HTTPException(status_code=502, detail=f"OpenAI API processing failed: {r.status_code}")
        return r.json()

@api_router.post("/parse-itinerary")
async def parse_itinerary(input: ParseItineraryInput, user: Any = Depends(verify_token)):
    api_key = os.environ.get("OPENAI_API_KEY")
    if not api_key:
        raise HTTPException(status_code=500, detail="OpenAI API key not configured on backend")
    
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }
    
    payload = {
        "model": "gpt-4o-mini",
        "messages": [
            {
                "role": "system",
                "content": (
                    "You are an expert travel assistant. Parse the following travel itinerary text and extract: "
                    "Days (with title, day number, and description), Hotels (name, location, price if any), "
                    "Activities, Transfers, Meals, and Notes. "
                    "Return a JSON object matching this schema: "
                    "{\n"
                    "  \"name\": \"Name of the tour/itinerary\",\n"
                    "  \"destination\": \"Main destination country/city\",\n"
                    "  \"days_count\": 7,\n"
                    "  \"days\": [\n"
                    "    {\n"
                    "      \"day\": 1,\n"
                    "      \"title\": \"Arrival in Paris\",\n"
                    "      \"description\": \"Morning VIP transfer...\",\n"
                    "      \"hotels\": [\"Hotel A\"],\n"
                    "      \"activities\": [\"Louvre private tour\"],\n"
                    "      \"transfers\": [\"Private chauffeur transfer\"],\n"
                    "      \"meals\": [\"Breakfast\", \"Dinner\"],\n"
                    "      \"notes\": \"Dress code is formal for dinner.\"\n"
                    "    }\n"
                    "  ],\n"
                    "  \"hotels\": [\n"
                    "    { \"name\": \"Hotel name\", \"location\": \"Hotel location\", \"price_per_night\": 420 }\n"
                    "  ],\n"
                    "  \"activities\": [\n"
                    "    { \"name\": \"Activity name\", \"price\": 75, \"description\": \"description\" }\n"
                    "  ]\n"
                    "}"
                )
            },
            {
                "role": "user",
                "content": input.text
            }
        ],
        "response_format": { "type": "json_object" }
    }
    
    try:
        result = await call_openai_with_retry(payload, headers)
        import json
        content = result["choices"][0]["message"]["content"]
        return json.loads(content)
    except Exception as e:
        logger.exception("parse itinerary failed")
        raise HTTPException(status_code=500, detail=str(e))


# ── PDF proxy ─────────────────────────────────────────────────────────────
# Forwards proposal export JSON to the Node Puppeteer service running on
# PDF_SERVICE_URL (default http://localhost:8002) and streams the resulting
# PDF back to the browser. Keeping this in FastAPI avoids any CORS or
# ingress-port issues because the React app only ever talks to /api/*.
PDF_SERVICE_URL = os.environ.get('PDF_SERVICE_URL', 'http://localhost:8002')

@api_router.get("/pdf/health")
async def pdf_health():
    try:
        async with httpx.AsyncClient(timeout=5.0) as c:
            r = await c.get(f"{PDF_SERVICE_URL}/health")
            return {"upstream_status": r.status_code, "upstream": r.json()}
    except Exception as e:
        return {"upstream_status": 0, "error": str(e)}

@api_router.post("/pdf/generate")
async def pdf_generate(request: PDFGenerateRequest, user: Any = Depends(verify_token)):
    payload = request.model_dump()
    user_id = user.get("id") or user.get("sub") or "unknown_user"
    try:
        async with httpx.AsyncClient(timeout=60.0) as c:
            r = await c.post(f"{PDF_SERVICE_URL}/generate", json=payload, headers={"X-User-ID": user_id})
        if r.status_code != 200:
            raise HTTPException(status_code=r.status_code, detail=r.text)
        proposal_name = (payload.get('proposal') or {}).get('name') or 'proposal'
        safe = ''.join(ch if ch.isalnum() or ch in ('.', '_', '-') else '-' for ch in proposal_name)
        return Response(
            content=r.content,
            media_type='application/pdf',
            headers={'Content-Disposition': f'attachment; filename="{safe}.pdf"'},
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.exception('pdf proxy failed')
        raise HTTPException(status_code=502, detail=f'PDF service unreachable: {e}')


# ── PPT generator ─────────────────────────────────────────────────────────
@api_router.post("/ppt/generate")
async def ppt_generate(request: PPTGenerateRequest, user: Any = Depends(verify_token)):
    payload = request.model_dump()
    try:
        from pptx import Presentation
        from io import BytesIO

        prs = Presentation()
        proposal = payload.get("proposal", {})
        proposal_name = proposal.get("name", "Voyanta Proposal")
        client_name = proposal.get("client_name", "")

        # Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = proposal_name
        subtitle.text = f"Prepared for: {client_name}" if client_name else "Exclusive Travel Proposal"

        # Overview Slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = "Itinerary Overview"
        tf = body_shape.text_frame
        
        items = payload.get("items", [])
        for item in items[:10]: # Limit to 10 for overview
            p = tf.add_paragraph()
            p.text = f"Day {item.get('day_number', '?')}: {item.get('title', item.get('kind', 'Item'))}"
            p.level = 0

        # Save to memory
        ppt_stream = BytesIO()
        prs.save(ppt_stream)
        ppt_stream.seek(0)
        
        safe = ''.join(ch if ch.isalnum() or ch in ('.', '_', '-') else '-' for ch in proposal_name)
        return Response(
            content=ppt_stream.read(),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{safe}.pptx"'},
        )
    except Exception as e:
        logger.exception("ppt generation failed")
        raise HTTPException(status_code=500, detail=f"PPT generation failed: {e}")

# Include the router in the main app
app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get('CORS_ORIGINS', 'http://localhost:5173').split(','),
    allow_methods=["*"],
    allow_headers=["*"],
)



if __name__ == "__main__":
    import uvicorn
    uvicorn.run("server:app", host="0.0.0.0", port=8001, reload=True)